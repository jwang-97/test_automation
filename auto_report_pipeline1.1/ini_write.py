#import os, sys
#import configParser

##From https://www.pythonf.cn/read/97851
#def ini_write(org_import_dir, config_file_path):
#    cf = configParser.ConfigParser()
#    cf.read(config_file_path)

#    #s = cf.sections()
#    #print 'section:', s

#    #o = cf.options("baseconf")
#    #print 'options:', o

#    #v = cf.items("baseconf")
#    #print 'db:', v
#    #dynamic_path = config_file_path + "\\Dynammic\\links"
#    #if exists(dynamic_path):
#    #    for root,filenames in os.walk(dynamic_path):
#    #        if name == "build":
#    #            dynamic_path = os.path.join(root, name)
#    #            cf.set("bit_base_directory", "bit1_directory", dynamic_path)
#    #            cf.set("bit_base_directory", "bit2_directory", dynamic_path)
#    g = os.walk(org_import_dir)
#    for path, dir_list, file_list in g:
#        for dir_name in dir_list:
#            if dir_name == "E616":
#                for sub_path, sub_dir_list, sub_file_list in os.walk(os.path.join(path, dir_name)):
#                    for file_name in sub_file_list:
#                        if file_name == "Wiley.des":
#                            cf.set("DESfile_E616_path", "cases_path1", os.path.join(sub_pathh, file_name))
#            elif dir_name == "stinger" or "axe" or "pdc" or "px" or "hyper" or "echo":
#                for sub_path, sub_dir_list, sub_file_list in os.walk(os.path.join(path, dir_name)):
#                    for file_name in sub_file_list:
#                        if file_name == ".des":
#                            type = dir_name.split('_')(-1)
#                            cf.set("DESfile_" + type + "_path", "cases_path1", os.path.join(sub_path, file_name))
#        #for file_name in file_list:
#        #    if name  == "pdc":
#        #        for root,filename in os.walk(name):

#            #if name == "stinger:
#            #elif name == "axe":
#            #elif name == "pdc":
#            #elif name == "px":
#            #elif name == "hyper":
#            #elif name == "E616":
#            #elif name == "echo":
            
    
#    config.write(open(config_file_path, "w"))
#    #db_host = cf.get("baseconf", "host")
#    #db_port = cf.getint("baseconf", "port")
#    #db_user = cf.get("baseconf", "user")
#    #db_pwd = cf.get("baseconf", "password")

#    #print db_host, db_port, db_user, db_pwd

#    #cf.set("baseconf", "db_pass", "123456")
#    #cf.write(open("config_file_path", "w"))
#if __name__ == "__main__":
#    ini_write(r'C:\Users\JWang294\Documents\projectfile\static-cases\Static', r'C:\Users\JWang294\Documents\projectfile\auto_test_config.ini')



import os, sys, re, logging, time, argparse, pptx, pptx.util, glob, re, csv, cv2, io
import configparser
import win32com.client as win32
import xlwings as xw
from pptx import Presentation
import matplotlib.pyplot as plt
import numpy as np

config_ini_dir = os.getcwd() + r'\auto_test_config.ini'
logging.basicConfig(format='%(asctime)s - %(pathname)s[line:%(lineno)d] - %(levelname)s: %(message)s',
                    level=logging.DEBUG)

def ini_write(org_import_dir1, org_import_dir2, config_file_path):
    cf = configparser.ConfigParser()
    #with open(config_file_path, 'rb') as f:
    #cf.read(config_file_path, encoding='utf-8', errors='ignore')
    #cf.read(config_file_path, "rb")
    cf.read(config_file_path, encoding='utf-8')
    org_import_dir = (org_import_dir1, org_import_dir2)
    i = 1
    while (i<3):
        g = os.walk(org_import_dir[i-1])
        for path, dir_list, file_list in g:
            for dir_name in dir_list:
                if re.search("E616", dir_name) != None:
                    for sub_path, sub_dir_list, sub_file_list in os.walk(os.path.join(path, dir_name)):
                        for file_name in sub_file_list:
                            if re.search("Wiley.des", file_name) != None:
                                cf.set("DESfile_E616_path", "cases_path"+ str(i), os.path.join(sub_path, file_name))
                                #cf.write(open(config_file_path, "w+", encoding='utf-8'))
                                break
                elif re.search("stinger", dir_name) or re.search("axe", dir_name) or re.search("pdc", dir_name) or re.search("px", dir_name) or  re.search("hyper", dir_name) or  re.search("echo", dir_name) != None:
                    for sub_path, sub_dir_list, sub_file_list in os.walk(os.path.join(path, dir_name)):
                        for file_name in sub_file_list:
                            if re.search('.des', file_name) != None:
                                type_name = dir_name.split('_')
                                cf.set("DESfile_" + type_name[1] + "_path", "cases_path"+ str(i), os.path.join(sub_path, file_name))
                                #cf.write(open(config_file_path, "w+", encoding='utf-8'))
                elif re.search("links", dir_name) != None:
                    for sub_path, sub_dir_list, sub_file_list in os.walk(os.path.join(path, dir_name)):
                        for sub_dir_name in sub_dir_list:
                            while re.match('build', sub_dir_name) != None:
                                #type_name = dir_name.split('_')
                                cf.set("bit_base_directory", "bit"+ str(i) +"_directory", os.path.join(sub_path, sub_dir_name))
                                #cf.write(open(config_file_path, "w+", encoding='utf-8'))
                                bln_build = True
                                break
                            if bln_build:
                                break
                        if bln_build:
                            break
                else:
                    logging.debug(dir_name + ' incompatible')
                    #print(dir_name + ' incompatible')
                    continue
        cf.write(open(config_file_path, "w+", encoding='utf-8'))
        i += 1

def orgpath_sequence(directory1, directory2):
    cf = configparser.ConfigParser()
    cf.read(config_ini_dir, encoding='utf-8')
    date1 = directory1.split('\\')[-1]
    date2 = directory2.split('\\')[-1]
    if date1.split("-")[1] > date2.split("-")[1]:
        #print(date1.split("-")[1])
        cf.set("profile_title", "title1", str(date2.split("-")[1])[0:8])
        cf.set("profile_title", "title2", str(date1.split("-")[1])[0:8])
        cf.write(open(config_ini_dir, "w+", encoding='utf-8'))
        return directory2, directory1
    else:
        cf.set("profile_title", "title1", str(date1.split("-")[1])[0:8])
        #print(str(date1.split("-")[1])[0:8])
        cf.set("profile_title", "title2", str(date2.split("-")[1])[0:8])
        cf.write(open(config_ini_dir, "w+", encoding='utf-8'))
        return directory1, directory2

def write_template_path_ini(tpl_path):
    cf = configparser.ConfigParser()
    cf.read(config_ini_dir, encoding='utf-8')
    if os.path.exists(tpl_path) == False:
        template_file = os.getcwd() + r'\60095160_RevA_CPC_Template_v4.04.pptx'
        cf.set("templatefile_path", "template_path", template_file)
        logging.debug('Invalid template path, default read template under the current path')
        logging.debug('ppttemplate-path:' + template_file)
    else:
        cf.set("templatefile_path", "template_path", tpl_path)
        logging.debug('ppttemplate-path:' + tpl_path)
    cf.write(open(config_ini_dir, "w+", encoding='utf-8'))
    return None
    
def call_SAM_macro():
    #try:
        #Launch Excel and Open Wrkbook
        xl=win32.Dispatch("Excel.Application")
        xl.Visible = True
        xl.Application.ScreenUpdating = False
        str = os.getcwd()
        xl.Workbooks.Open(Filename = str + "\\test_SAM_v11.78.xlsm")
        time.sleep(5)
        #xl.ThisWorkbook.Activate
        # xl.Workbooks("Sheet1").Activate
        #xl.Application.ScreenUpdating = True
        logging.debug('Open SAMfile success')
        #Run Macro
        xl.Application.Run('test_SAM_v11.78.xlsm!Module1.Auto_test')
        logging.debug('Auto_test')
        # xl.Application.ScreenUpdating = True
        xl.ScreenUpdating = True
        # xl.Application.DisplayAlerts = False
        xl.DisplayAlerts = False
        xl.ActiveWorkbook.Close(SaveChanges:=False)
        xl.Quit
      
def call_Multi_macro():
    xl=win32.Dispatch("Excel.Application")
    xl.Visible = True
    xl.Application.ScreenUpdating = False
    str = os.getcwd()
    xl.Workbooks.Open(Filename = str + "\\Multi_Bit_Compare_v2.52.xlsm")
    time.sleep(5)
    logging.debug('Open Multifile success')
    #Run Macro
    xl.Application.Run('Multi_Bit_Compare_v2.52.xlsm!Module1.Auto_Dynamic')
    logging.debug('Auto_Dynamic')
    xl.Application.ScreenUpdating = True
    xl.ScreenUpdating = True
    # xl.Application.DisplayAlerts = False
    xl.DisplayAlerts = False
    xl.ActiveWorkbook.Close(SaveChanges:=False)
    xl.Quit
    
def read_doc_rop(org_dir1):
    i = 0
    doc_list = []
    rop_list = []
    for pic_dir in glob.glob(org_dir1 + r'\Contact\contact_script_test\summary_doc'):
        # rop_data = pd.read_csv(pic_dir, sep = '\t',encoding = 'utf-8', errors="ignore")
        # rop_data.head()
        with open(pic_dir, "r", encoding="gbk", errors="ignore") as rop_data:
            csv_reader = csv.reader(rop_data)
            for row in csv_reader:
                #save rop as a list
                if i>=3:
                    # rop_list[i-3] = list(str(row).split())[3]
                    doc_list = str(row).split()
                    print(isinstance(doc_list ,list))
                    rop_list.append(doc_list[4])
                    # print(rop_list[i-3])
                    # print(rop_list)
                i += 1
        break
    return rop_list
    
def contact_pic_hooks(org_import_dir1, org_import_dir2, config_file_path):
    cf = configparser.ConfigParser()
    cf.read(config_file_path, encoding='utf-8')
    template_path = cf.get("merge_template", "dirct")
    prs = Presentation(template_path)
    org_import_dir = (org_import_dir1, org_import_dir2)
    contact_pic_path = ['','']
    contact_pic = ['','']
    i = 0
    while i <= 1:
        # org_path = os.walk(org_import_dir[i+1])
        # for path, dir_list, file_list in org_path:
            # for dir_name in dir_list:
                # if re.search("\Contact\contact_script_test\utility_files", dir_name) != None:
                    # logging.debug('invalid contact path:' + org_path + ", lacks essential utility_files")
        for pic_dir in glob.glob(org_import_dir[i] + r'\Contact\contact_script_test\utility_files*'):
                    # cf.set("Contact_pic_path", "pic_path"+ str(i), os.path.join(path, dir_name))
            contact_pic_path[i] = pic_dir
            break
        i += 1
    pic_left = [int(prs.slide_width*0.08), int(prs.slide_width*0.58)]
    pic_top = int(prs.slide_height*0.3)
    pic_width = int(prs.slide_width*0.4)
    pic_height = int(prs.slide_height*0.6)
    ipath = 0
    time_title1 = cf.get("profile_title", "title1")
    time_title2 = cf.get("profile_title", "title2")
    #get rop number
    rop_num_list = read_doc_rop(org_import_dir1)
    # text_array = ("ROP Controlled: " + "50" + " ft/hr", "Build"+ time_title1, "Build" + time_title2)
    # while ipath <= 1:
    for contact_pic_path1 in glob.glob(contact_pic_path[ipath]+'/*.jpg'):
        if re.search('Blade', contact_pic_path1) == None:
            pic_type = contact_pic_path1.split('\\')[-1]
            #get number of pic to match rop
            bits_num = list(re.findall(r"\d+" ,str(pic_type.split('_')[0])))[-1]
            text_array = ("ROP Controlled: " + str(rop_num_list[int(bits_num) - 1]) + " ft/hr", "Build"+ time_title1, "Build" + time_title2)
            for contact_pic_path2 in glob.glob(contact_pic_path[ipath+1]+'/' + pic_type + '*'):
                # contact_pic[1] = contact_pic_path2
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                # slide.shapes.title.textframe.textrange.text = "test contact slide"
                slide.shapes.title.text = "Contact Analysis Comparison - Blade Top Contact"
                # pic2 = slide.shapes.add_picture()
                img1 = cv2.imread(contact_pic_path[0])
                img2 = cv2.imread(contact_pic_path[1])
                pic1 = slide.shapes.add_picture(contact_pic_path1, pic_left[0], pic_top, pic_width, pic_height)
                pic2 = slide.shapes.add_picture(contact_pic_path2, pic_left[1], pic_top, pic_width, pic_height)
                # print (contact_pic_path1)
                i = 0
                while i <= 2:
                    if i == 0:
                        tb = slide.shapes.add_textbox(300000, 700000, 1, 1)
                    elif i == 1:
                        tb = slide.shapes.add_textbox(1000000, 1600000, 1, 1)
                    else:
                        tb = slide.shapes.add_textbox(8000000, 1600000, 1, 1)
                    p1 =  tb.text_frame.add_paragraph()
                    p1.text = text_array[i]
                    p1.font.size = pptx.util.Pt(25)
                    i+=1

    prs.save(template_path)

def read_summary_doc(org_dir, list_type):
    i = 0
    iCase = 0
    # doc_list = []
    axis_num_list = []
    # Connect zero point
    axis_num_list.append(0)
    for pic_dir in glob.glob(org_dir + r'\Contact\contact_script_test\summary_doc'):
        with open(pic_dir, "r", encoding="gbk", errors="ignore") as rop_data:
            csv_reader = csv.reader(rop_data)
            for row in csv_reader:
                doc_list = str(row).split()
                #save rop/wob... as a list
                if i >= 2 and i < 3:
                    # doc_list = str(row).split()
                    for type in doc_list:
                        #match type to choose which line to return
                        if (type.find(list_type)) != -1:
                            iCase = doc_list.index(type)
                            break
                #add right num_line to axis_num_list
                elif i > 3:
                    axis_num_list.append(float(doc_list[iCase]))
                i += 1
        break
    return axis_num_list

def generate_contact_plots(org_dir1, org_dir2, config_path):
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    template_path = cf.get("merge_template", "dirct")
    prs = Presentation(template_path)
    org_import_dir = (org_dir1, org_dir2)
    time_title1 = cf.get("profile_title", "title1")
    time_title2 = cf.get("profile_title", "title2")
    doc_type_list = ("doc", "tq", "rop")
    pic_left = int(prs.slide_width*0.1)
    pic_top = int(prs.slide_height*0.08)
    pic_width = int(prs.slide_width*0.6)
    pic_height = int(prs.slide_height*0.8)
    for type in doc_type_list:
        # line 1 points
        x1 = read_summary_doc(org_dir1, "wob")
        y1 = read_summary_doc(org_dir1, type)
        # plotting the line 1 points 
        plt.plot(x1, y1, label = "build" + time_title1, color='c')
        # line 2 points
        x2 = read_summary_doc(org_dir2, "wob")
        y2 = read_summary_doc(org_dir2, type)
        # plotting the line 2 points 
        plt.plot(x2, y2, label = "build" + time_title2, color='m')
        # add Horizontal line
        plt.grid(axis="y")
        # naming the x&y axis
        if type == "doc":
            y_axis_label = "Depth of Cut (in/rev)"
            plot_title = "Depth of Cut-Wellington Shale @ 3ksi"
            # add Horizontal line
            plt.axhline(y2[3], color='orange', lw=2, alpha=0.7)
            # fill with color
            plt.axhspan(0, y2[3], alpha=0.5, color='y')
            # plt.fill_between(x = , y = y2[3], y = 0, step="pre", color = "g", alpha = 0.3)
            plt.ylim(0,0.6)
        elif type == "tq":
            y_axis_label = "Torque (ft-lbf)"
            plot_title = "Torque-Wellington Shale @ 3ksi" 
            plt.axvline(x2[4], color='gold', lw=2, alpha=0.7)
            plt.ylim(0,10000)
        else:
            y_axis_label = "ROP (ft/hr)"
            plot_title = "ROP-Wellington Shale @ 3ksi" 
            plt.axvline(x2[4], color='orange', lw=2, alpha=0.7)
            plt.ylim(0,450)
            # plt.axis([10000,60000, 50,450])
        plt.xlim(0,)
        plt.xlabel('WOB(lbf)')
        plt.ylabel(y_axis_label)
        # giving a title to my graph
        plt.title(plot_title)
        plt.legend()
        buffer_png = io.BytesIO()
        plt.savefig(buffer_png, format = "png")
        # plt.show()
        # # function to show the plot
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        pic1 = slide.shapes.add_picture(buffer_png, pic_left[0], pic_top, pic_width, pic_height)
        plt.close()
        buffer_png.close()
    prs.save(template_path)

if __name__ == "__main__":
    cf = configparser.ConfigParser()
    # config_ini_dir = os.getcwd() + r'\auto_test_config.ini'
    cf.read(config_ini_dir, encoding='utf-8')
    logging.debug('config file path:' + config_ini_dir)
    # org_result_path1, org_result_path2, template_path = input("enter cases-results path1 & cases-result path2 & ppt_template-path separated with ',':").split(',')
    parser = argparse.ArgumentParser(description='input path argparse')
    parser.add_argument('--path1', '-f', help='path1 property, non-essential parameters, have default values', default= r"U:\result-backup\test-2021\test-20210416173550-lbo-0416cpc")
    parser.add_argument('--path2', '-s', help='path2 property, non-essential parameters, have default values', default= r"U:\result-backup\test-2021\test-20210702174422-lbo-0630cpc")
    parser.add_argument('--template_path', '-t', help='template_path property, non-essential parameters, have default values', default="")
    args = parser.parse_args()
    org_result_path1 = args.path1
    org_result_path2 = args.path2
    template_path = args.template_path
    if os.path.exists(org_result_path1) == False:
        org_result_path1 = cf.get("original_path", "org_path1")
        logging.debug('Invalid path1, default read original_path in auto_test_config.ini')

    if os.path.exists(org_result_path2) == False:
        org_result_path2 = cf.get("original_path", "org_path2")
        logging.debug('Invalid path2, default read original_path in auto_test_config.ini')

    logging.debug('path1:' + org_result_path1)
    logging.debug('path2:' + org_result_path2)
    write_template_path_ini(template_path)
    # write several types of path in config.ini 
    config_dir = orgpath_sequence(org_result_path1, org_result_path2)
    ini_write(config_dir[0], config_dir[1], config_ini_dir)
    
    call_SAM_macro()
    generate_contact_plots(org_result_path1, org_result_path2, config_ini_dir)
    contact_pic_hooks(org_result_path1, org_result_path2, config_ini_dir)
    call_Multi_macro()
    
    # run_command = "start " + os.getcwd() + r'\test_SAM_v11.78.xlsm'
    # os.system(run_command)
    #ini_write(r'\\bgc-nas002\ideas_engine_folder\Share\Drilling_software_development\Work\Work-2021\release_testing_for_cpc\testing-cases', r'C:\Users\JWang294\Documents\projectfile\auto_test_config.ini')
    #config_sequence(r'\\bgc-nas002\ideas_engine_folder\Share\Drilling_software_development\Work\Work-2021\release_testing_for_cpc\testing-cases', r'\\bgc-nas002\ideas_engine_folder\Share\Drilling_software_development\Work\Work-2021\release_testing_for_cpc\testing-cases')
    # C:\Users\JWang294\Documents\projectfile\test_config_cases\test-20210416173550-lbo-0416cpc
    # C:\Users\JWang294\Documents\projectfile\test_config_cases\test-20210702174422-lbo-0630cpc