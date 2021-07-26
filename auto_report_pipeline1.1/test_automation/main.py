import os, re, logging, time, argparse, pptx, pptx.util, csv, io
import configparser
import win32com.client as win32
# import xlwings as xw
from pptx import Presentation
import matplotlib.pyplot as plt
# import numpy as np
from pptx.util import Cm
import pandas as pd

logging.basicConfig(format='%(asctime)s - %(pathname)s[line:%(lineno)d] - %(levelname)s: %(message)s',
                    level=logging.DEBUG)

def ini_write(org_import_dir1, org_import_dir2, config_file_path):
    cf = configparser.ConfigParser()
    #with open(config_file_path, 'rb') as f:
    #cf.read(config_file_path, encoding='utf-8', errors='ignore')
    #cf.read(config_file_path, "rb")
    cf.read(config_file_path, encoding='utf-8')
    org_import_dir = (org_import_dir1, org_import_dir2)
    i = 1
    while (i<3):
        g = os.walk(org_import_dir[i-1])
        #write file_path in ini-config file
        for path, dir_list, file_list in g:
            for dir_name in dir_list:
                #write E616 file_path
                if re.search("E616", dir_name) != None:
                    for sub_path, sub_dir_list, sub_file_list in os.walk(os.path.join(path, dir_name)):
                        for file_name in sub_file_list:
                            if re.search("Wiley.des", file_name) != None:
                                section_name = "DESfile_" + "E616" + "_path"
                                if not cf.has_section(section_name):
                                    cf.add_section(section_name)
                                cf.set(section_name, "cases_path"+ str(i), os.path.join(sub_path, file_name))
                                #cf.write(open(config_file_path, "w+", encoding='utf-8'))
                                break
                #write stinger file_path
                elif re.search("stinger", dir_name) or re.search("axe", dir_name) or re.search("pdc", dir_name) or re.search("px", dir_name) or  re.search("hyper", dir_name) or  re.search("echo", dir_name) != None:
                    for sub_path, sub_dir_list, sub_file_list in os.walk(os.path.join(path, dir_name)):
                        for file_name in sub_file_list:
                            if re.search('.des', file_name) != None:
                                type_name = dir_name.split('_')
                                section_name = "DESfile_" + type_name[1] + "_path"
                                if not cf.has_section(section_name):
                                    cf.add_section(section_name)
                                cf.set(section_name, "cases_path"+ str(i), os.path.join(sub_path, file_name))
                                #cf.write(open(config_file_path, "w+", encoding='utf-8'))
                #write dynamic file_path
                elif re.search("links", dir_name) != None:
                    for sub_path, sub_dir_list, sub_file_list in os.walk(os.path.join(path, dir_name)):
                        for sub_dir_name in sub_dir_list:
                            while re.match('build', sub_dir_name) != None:
                                #type_name = dir_name.split('_')
                                section_name = "bit_base_directory"
                                if not cf.has_section(section_name):
                                    cf.add_section(section_name)
                                cf.set(section_name, "bit"+ str(i) +"_directory", os.path.join(sub_path, sub_dir_name))
                                #cf.write(open(config_file_path, "w+", encoding='utf-8'))
                                bln_build = True
                                break
                            if bln_build:
                                break
                        if bln_build:
                            break
                #write contact file_path
                elif re.search("contact_script_test", dir_name) != None:
                    for sub_path, sub_dir_list, sub_file_list in os.walk(os.path.join(path, dir_name)):
                        for file_name in sub_file_list:
                            if re.search('summary_doc', file_name) != None:
                                # type_name = dir_name.split('_')
                                section_name = "contact_summary_doc_dir"
                                if not cf.has_section(section_name):
                                    cf.add_section(section_name)
                                cf.set(section_name, "cases_path"+ str(i), os.path.join(sub_path, file_name))
                        for sub_dir_name in sub_dir_list:
                            if re.search('utility_files', sub_dir_name) != None:
                                # type_name = dir_name.split('_')
                                section_name = "contact_pic_dir"
                                if not cf.has_section(section_name):
                                    cf.add_section(section_name)
                                cf.set(section_name, "cases_path"+ str(i), os.path.join(sub_path, sub_dir_name))
                #write other-type static file_path 
                elif re.search("6", dir_name) != None and re.search("_", dir_name) != None and len(dir_name.split("_")) == 2:
                    for sub_path, sub_dir_list, sub_file_list in os.walk(os.path.join(path, dir_name)):
                        for file_name in sub_file_list:
                            if re.search('.des', file_name) != None:
                                type_name = dir_name.split('_')
                                section_name = "DESfile_" + type_name[1] + "_path"
                                if not cf.has_section(section_name):
                                    cf.add_section(section_name)
                                #     cf.set(section_name, "cases_path"+ str(i), os.path.join(sub_path, file_name))
                                # else:
                                cf.set(section_name, "cases_path"+ str(i), os.path.join(sub_path, file_name))
                else:
                    logging.debug(dir_name + ' incompatible')
                    #print(dir_name + ' incompatible')
                    continue
        cf.write(open(config_file_path, "w+", encoding='utf-8'))
        i += 1

#write release_time into config.ini
def orgpath_sequence(directory1, directory2 ,config_path):
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    date1 = directory1.split('\\')[-1]
    date2 = directory2.split('\\')[-1]
    if date1.split("-")[1] > date2.split("-")[1]:
        #print(date1.split("-")[1])
        cf.set("profile_title", "title1", str(date2.split("-")[1])[0:8])
        cf.set("profile_title", "title2", str(date1.split("-")[1])[0:8])
        cf.write(open(config_path, "w+", encoding='utf-8'))
        return directory2, directory1
    else:
        cf.set("profile_title", "title1", str(date1.split("-")[1])[0:8])
        #print(str(date1.split("-")[1])[0:8])
        cf.set("profile_title", "title2", str(date2.split("-")[1])[0:8])
        cf.write(open(config_path, "w+", encoding='utf-8'))
        return directory1, directory2

#verification ppt template path
def write_template_path_ini(tpl_path, config_path):
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    if os.path.exists(tpl_path) == False:
        path = os.getcwd()
        # parent = os.path.abspath(os.path.join(path, os.pardir))
        template_file = path + r'\60095160_RevA_CPC_Template_v4.04.pptx'
        cf.set("templatefile_path", "template_path", template_file)
        logging.debug('Invalid template path, default read template under the current path')
        logging.debug('ppttemplate-path:' + template_file)
    else:
        cf.set("templatefile_path", "template_path", tpl_path)
        logging.debug('ppttemplate-path:' + tpl_path)
    cf.write(open(config_path, "w+", encoding='utf-8'))
    return None

#call static macros in SAM
def call_SAM_macro():
    #try:
        #Launch Excel and Open Wrkbook
        xl=win32.Dispatch("Excel.Application")
        xl.Visible = True
        xl.Application.ScreenUpdating = False
        str = os.getcwd()
        xl.Workbooks.Open(Filename = str + r"\test_automation\SAM_v11.78.xlsm")
        time.sleep(5)
        #xl.ThisWorkbook.Activate
        # xl.Workbooks("Sheet1").Activate
        #xl.Application.ScreenUpdating = True
        logging.debug('Open SAMfile success')
        #Run Macro
        xl.Application.Run('SAM_v11.78.xlsm!Module1.Auto_test')
        logging.debug('Auto_test')
        # xl.Application.ScreenUpdating = True
        xl.ScreenUpdating = True
        # xl.Application.DisplayAlerts = False
        xl.DisplayAlerts = False
        # xl.ActiveWorkbook.Close(SaveChanges:=False)
        xl.ActiveWorkbook.Close()
        xl.Quit
      
#call dynamic macros in Multi
def call_Multi_macro():
    xl=win32.Dispatch("Excel.Application")
    xl.Visible = True
    xl.Application.ScreenUpdating = False
    str = os.getcwd()
    xl.Workbooks.Open(Filename = str + r"\test_automation\Multi_Bit_Compare_v2.52.xlsm")
    time.sleep(5)
    logging.debug('Open Multifile success')
    #Run Macro
    xl.Application.Run('Multi_Bit_Compare_v2.52.xlsm!Module1.Auto_Dynamic')
    logging.debug('Auto_Dynamic')
    xl.Application.ScreenUpdating = True
    xl.ScreenUpdating = True
    # xl.Application.DisplayAlerts = False
    xl.DisplayAlerts = False
    # xl.ActiveWorkbook.Close(SaveChanges:=False)
    xl.ActiveWorkbook.Close()
    xl.Quit
    
#get contact rop-list in summary_doc
def read_doc_rop(iCase, config_path):
    i = 0
    doc_list = []
    rop_list = []
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    pic_dir = cf.get("contact_summary_doc_dir", "cases_path" + str(iCase))
    # for pic_dir in glob.glob(org_dir1 + r'**\Contact\contact_script_test\summary_doc'):
        # rop_data = pd.read_csv(pic_dir, sep = '\t',encoding = 'utf-8', errors="ignore")
        # rop_data.head()
    with open(pic_dir, "r", encoding="gbk", errors="ignore") as rop_data:
        csv_reader = csv.reader(rop_data)
        for row in csv_reader:
            #save rop as a list
            if i>=3:
                # rop_list[i-3] = list(str(row).split())[3]
                doc_list = str(row).split()
                # print(isinstance(doc_list ,list))
                rop_list.append(doc_list[4])
                # print(rop_list[i-3])
                # print(rop_list)
            i += 1
            # break
    return rop_list

#get pictures under the path
def get_img_file(path_name):
    imagelist = []
    for parent, dirnames, filenames in os.walk(path_name):
        for filename in filenames:
            if filename.lower().endswith(('.bmp', '.dib', '.png', '.jpg', '.jpeg', '.pbm', '.pgm', '.ppm', '.tif', '.tiff')):
                imagelist.append(os.path.join(parent, filename))
        return imagelist

#get contact pictures
# def contact_pic_hooks(org_import_dir1, org_import_dir2, config_file_path):
def contact_pic_hooks(config_file_path):
    cf = configparser.ConfigParser()
    cf.read(config_file_path, encoding='utf-8')
    template_path = cf.get("merge_template", "dirct")
    prs = Presentation(template_path)
    # org_import_dir = (org_import_dir1, org_import_dir2)
    contact_pic_path = ['','']
    # contact_pic = ['','']
    i = 0
    while i <= 1:
        #get list of contact_pictures' path
        contact_pic_path[i] = get_img_file(cf.get("contact_pic_dir", "cases_path" + str(i+1)))
        # for pic_dir in glob.glob(org_import_dir[i] + r'*\Contact\contact_script_test\utility_files*'):
        #     contact_pic_path[i] = pic_dir
        #     break
        i += 1
    pic_left = [int(prs.slide_width*0.08), int(prs.slide_width*0.58)]
    pic_top = int(prs.slide_height*0.3)
    pic_width = int(prs.slide_width*0.4)
    pic_height = int(prs.slide_height*0.6)
    ipath = 0
    time_title1 = cf.get("profile_title", "title1")
    time_title2 = cf.get("profile_title", "title2")
    #get rop number
    rop_num_list = read_doc_rop(1, config_file_path)
    inum = 0
    while ipath < len(contact_pic_path[0]):
    # for contact_pic1 in contact_pic_path[0]:
        # generate plots in order
        # if re.search('Blade', contact_pic1) == None and re.findall(r"\d+", contact_pic1)[-1] == str(inum + 1):
        contact_pic1 = list(contact_pic_path[0])[ipath]
        if re.search('Blade', contact_pic1) == None and re.findall(r"\d+", contact_pic1)[-1] == str(inum + 1):
            pic_type = contact_pic1.split('\\')[-1]
            text_array = ("ROP Controlled: " + str(rop_num_list[inum]) + " ft/hr", "Build"+ time_title1, "Build" + time_title2)
            for contact_pic2 in contact_pic_path[1]:
                if re.search(pic_type, contact_pic2) != None and re.search('Blade', contact_pic2) == None:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = "Contact Analysis Comparison - Blade Top Contact"
                    # img1 = cv2.imread(contact_pic_path[0])
                    # img2 = cv2.imread(contact_pic_path[1])
                    pic1 = slide.shapes.add_picture(contact_pic1, pic_left[0], pic_top, pic_width, pic_height)
                    pic2 = slide.shapes.add_picture(contact_pic2, pic_left[1], pic_top, pic_width, pic_height)
                    # print (contact_pic_path1)
                    i = 0
                    while i <= 2:
                        if i == 0:
                            tb = slide.shapes.add_textbox(300000, 700000, 1, 1)
                        elif i == 1:
                            tb = slide.shapes.add_textbox(1000000, 1600000, 1, 1)
                        else:
                            tb = slide.shapes.add_textbox(8000000, 1600000, 1, 1)
                        p1 =  tb.text_frame.add_paragraph()
                        p1.text = text_array[i]
                        p1.font.size = pptx.util.Pt(25)
                        i+=1
            #redefine ipath,refreash the loop
            ipath = -1
            inum += 1
        ipath += 1
    prs.save(template_path)

#get contact-list in summary_doc
def read_summary_doc(case_num, list_type, config_path):
    i = 0
    iCase = 0
    # doc_list = []
    axis_num_list = []
    # Connect zero point
    axis_num_list.append(0)
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    doc_dir = cf.get("contact_summary_doc_dir", "cases_path" + str(case_num))
    # for pic_dir in glob.glob(org_dir + r'\**\contact_script_test\summary_doc'):
    with open(doc_dir, "r", encoding="gbk", errors="ignore") as rop_data:
        csv_reader = csv.reader(rop_data)
        for row in csv_reader:
            doc_list = str(row).split()
            #save rop/wob... as a list
            if i == 2:
                # doc_list = str(row).split()
                for type in doc_list:
                    #match type to choose which line to return
                    if (type.find(list_type)) != -1:
                        iCase = doc_list.index(type)
                        break
                #add right num_line to axis_num_list
            elif i > 2:
                axis_num_list.append(float(doc_list[iCase]))
            i += 1
    return axis_num_list

#generate contact plots
def generate_contact_plots(org_dir1, org_dir2, config_path):
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    template_path = cf.get("merge_template", "dirct")
    prs = Presentation(template_path)
    org_import_dir = (org_dir1, org_dir2)
    time_title1 = cf.get("profile_title", "title1")
    time_title2 = cf.get("profile_title", "title2")
    doc_type_list = ("doc", "tq", "rop")
    pic_left = int(prs.slide_width*0.13)
    pic_top = int(prs.slide_height*0.1)
    pic_width = int(prs.slide_width*0.6)
    pic_height = int(prs.slide_height*0.8)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Contact Analysis Inputs"
    iCase = 0
    Input_name_list = str(cf.get("Contact_Analysis_Inputs", "Inputs_name")).split('|')
    Input_data_list = str(cf.get("Contact_Analysis_Inputs", "Inputs_data")).split('|')
    while iCase < 2: 
        if iCase == 0:
            tb = slide.shapes.add_textbox(950000, 1600000, 1, 1)
        else:
            tb = slide.shapes.add_textbox(8000000, 1600000, 1, 1)
        p1 = tb.text_frame.add_paragraph()
        if iCase == 0:
            list = Input_name_list
        else:
            list = Input_data_list
        f = io.StringIO()
        for name in list:
            f.write(name + '\n')
        p1.text = f.getvalue()
        f.close()
        p1.font.size = pptx.util.Pt(29)
        iCase += 1
    for type in doc_type_list:
        # line 1 points
        x1 = read_summary_doc(1, "wob", config_path)
        y1 = read_summary_doc(1, type, config_path)
        # plotting the line 1 points 
        plt.plot(x1, y1, label = "build" + time_title1, color='c')
        # line 2 points
        x2 = read_summary_doc(2, "wob", config_path)
        y2 = read_summary_doc(2, type, config_path)
        # plotting the line 2 points 
        plt.plot(x2, y2, label = "build" + time_title2, color='m')
        # add Horizontal line
        plt.grid(axis="y")
        # naming the x&y axis
        if type == "doc":
            y_axis_label = "Depth of Cut (in/rev)"
            plot_title = "Depth of Cut-Wellington Shale @ 3ksi"
            # add Horizontal line
            plt.axhline(y2[4], color='orange', lw=2, alpha=0.7)
            # fill with color
            plt.axhspan(0, y2[4], alpha=0.5, color='y')
            # plt.fill_between(x = , y = y2[3], y = 0, step="pre", color = "g", alpha = 0.3)
            plt.ylim(0,0.6)
        elif type == "tq":
            y_axis_label = "Torque (ft-lbf)"
            plot_title = "Torque-Wellington Shale @ 3ksi" 
            plt.axvline(x2[4], color='orange', lw=2, alpha=0.7)
            plt.ylim(0,10000)
        else:
            y_axis_label = "ROP (ft/hr)"
            plot_title = "ROP-Wellington Shale @ 3ksi" 
            plt.axvline(x2[4], color='orange', lw=2, alpha=0.7)
            plt.ylim(0,450)
            # plt.axis([10000,60000, 50,450])
        plt.xlim(0,)
        plt.xlabel('WOB(lbf)')
        plt.ylabel(y_axis_label)
        # giving a title to my graph
        plt.title(plot_title)
        plt.legend()
        buffer_png = io.BytesIO()
        plt.savefig(buffer_png, format = "png")
        # plt.show()
        # # function to show the plot
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        pic1 = slide.shapes.add_picture(buffer_png, pic_left, pic_top, pic_width, pic_height)
        plt.close()
        buffer_png.close()
    prs.save(template_path)

def function_check_sheet(config_path):
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    template_path = cf.get("merge_template", "dirct")
    prs = Presentation(template_path)
    df = pd.DataFrame(data={'Tool/Program': None,
                            'Pass/Fail': None,
                        },
                    index = ['I-Design (idd) / Scale', 'I-Design (idd) / Modify Cutter/Position',
                             'IDEAS Static (ids) / Normal Functionality', 'IDEAS Static (ids) / Multi-Core processing',
                             'IDEAS Dynamic (id3) / Normal Functionality', 'IDEAS Dynamic (id3) / Multi-Core processing',
                             'IDEAS Autoforce (ida) / Force Balance Optimization', 'IDEAS3 Utility (id3u)',
                             'IDEAS Analysis Request (iAR) / Create Projects (ideasbuild'+cf.get("profile_title", "title2")+')', 'Bit Design Toolbox / IAR Project Renamer', 
                             'Bit Design Toolbox / Rollback IDEAS Project', 'IDM / Create/Run Projectgroup',
                             'SAM', 'Multi-Bit']
                    )
    slide = prs.slides[3]
    # slide = prs.slides.add_slide(prs.slide_layouts[1])
    # slide.shapes.title.text = "Contact Analysis Comparison - Blade Top Contact"
    rows = 15
    columns = 2
    left = Cm(3)   
    top = Cm(4) 
    width = Cm(32) 
    height = Cm(0.1)
    table = slide.shapes.add_table(rows=rows,
                            cols=columns,
                            left=left,
                            top=top,
                            width=width,
                            height=height
                            )
    table = table.table
    # colunm_name
    header = df.columns
    for i, h in enumerate(header):
        cell = table.cell(0, i)
        cell.text = h
    # row_name
    row_names = df.index
    for i, r in enumerate(row_names):
        cell = table.cell(i+1, 0)
        cell.text = r
    prs.save(template_path)

# creat a list of static types in config.ini, those types would be shown in ppt.
def cases_types_config(config_file_path):
    cf = configparser.ConfigParser()
    cf.read(config_file_path, encoding='utf-8')
    type_list = []
    #find all sections in config.ini
    cases_types = cf.sections()
    for type in cases_types:
        if re.search('DESfile', type) != None and re.search('path', type) != None:
            #confirm the section is not empty 
            if cf.get(type, "cases_path1") != None and cf.get(type, "cases_path1") != None:
                type_list.append(type)
    str_type = ' '.join(str(type) for type in type_list)
    cf.set("static_case_type", "type_list", str_type)
    cf.write(open(config_file_path, "w+", encoding='utf-8'))

#Clean the static value in the config.ini
def clean_config_file(config_file_path):
    cf = configparser.ConfigParser()
    cf.read(config_file_path, encoding='utf-8')
    #find all sections in config.ini
    cases_types = cf.sections()
    # rsserve_list = ['templatefile', 'original', 'GEMs_Ref_Info', 'merge_template', 'Contact_Analysis_Inputs']
    for type in cases_types:
        if re.search('path', type) != None and re.search('DESfile', type) != None:
            cf.set(type, "cases_path1", "")
            cf.set(type, "cases_path2", "")
    cf.write(open(config_file_path, "w+", encoding='utf-8'))

if __name__ == "__main__":
    path = os.getcwd()
    # parent = os.path.abspath(os.path.join(path, os.pardir))
    config_ini_dir = path + r'\docs\auto_test_config.ini'
    cf = configparser.ConfigParser()
    # config_ini_dir = os.getcwd() + r'\auto_test_config.ini'
    cf.read(config_ini_dir, encoding='utf-8')
    logging.debug('config file path:' + config_ini_dir)
    # org_result_path1, org_result_path2, template_path = input("enter cases-results path1 & cases-result path2 & ppt_template-path separated with ',':").split(',')
    parser = argparse.ArgumentParser(description='input path argparse')
    parser.add_argument('--path1', '-f', help='path1 property, non-essential parameters, have default values', default= r"U:\result-backup\test-2021\test-20210416173550-lbo-0416cpc")
    parser.add_argument('--path2', '-s', help='path2 property, non-essential parameters, have default values', default= r"U:\result-backup\test-2021\test-20210702174422-lbo-0630cpc")
    parser.add_argument('--template_path', '-t', help='template_path property, non-essential parameters, have default values', default="")
    args = parser.parse_args()
    org_result_path1 = args.path1
    org_result_path2 = args.path2
    template_path = args.template_path
    if os.path.exists(org_result_path1) == False:
        org_result_path1 = cf.get("original_path", "org_path1")
        logging.debug('Invalid path1, default read original_path in auto_test_config.ini')

    if os.path.exists(org_result_path2) == False:
        org_result_path2 = cf.get("original_path", "org_path2")
        logging.debug('Invalid path2, default read original_path in auto_test_config.ini')

    logging.debug('path1:' + org_result_path1)
    logging.debug('path2:' + org_result_path2)
    clean_config_file(config_ini_dir)
    write_template_path_ini(template_path, config_ini_dir)
    # write several types of path in config.ini 
    config_dir = orgpath_sequence(org_result_path1, org_result_path2, config_ini_dir)
    ini_write(config_dir[0], config_dir[1], config_ini_dir)
    cases_types_config(config_ini_dir)
    call_SAM_macro()
    function_check_sheet(config_ini_dir)
    generate_contact_plots(config_dir[0], config_dir[1], config_ini_dir)
    contact_pic_hooks(config_ini_dir)
    call_Multi_macro()
    

    #ini_write(r'\\bgc-nas002\ideas_engine_folder\Share\Drilling_software_development\Work\Work-2021\release_testing_for_cpc\testing-cases', r'C:\Users\JWang294\Documents\projectfile\auto_test_config.ini')
    #config_sequence(r'\\bgc-nas002\ideas_engine_folder\Share\Drilling_software_development\Work\Work-2021\release_testing_for_cpc\testing-cases', r'\\bgc-nas002\ideas_engine_folder\Share\Drilling_software_development\Work\Work-2021\release_testing_for_cpc\testing-cases')
    # C:\Users\JWang294\Documents\projectfile\test_config_cases\test-20210416173550-lbo-0416cpc
    # C:\Users\JWang294\Documents\projectfile\test_config_cases\test-20210702174422-lbo-0630cpc