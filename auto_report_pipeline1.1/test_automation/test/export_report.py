import os, re, logging, time, argparse, pptx, pptx.util, csv, io, datetime
from pickle import TRUE
import configparser
import win32com.client as win32
# import xlwings as xw
from pptx import Presentation
import matplotlib.pyplot as plt
# import numpy as np
from pptx.util import Cm
import pandas as pd
import numpy as np
from pandas.plotting import table
from pylab import mpl
import matplotlib.pyplot as plt
from PIL import ImageGrab
from io import StringIO, BytesIO
from PIL import Image

#clear the csv file in the dir
def clear_csv(config_ini_dir):
    parent_path = config_ini_dir.replace(config_ini_dir.split("\\")[-1], '')
    csvlist = []
    for parent, dirnames, filenames in os.walk(parent_path):
        for filename in filenames:
            if filename.lower().endswith(('.csv')):
                csvlist.append(os.path.join(parent, filename))
    for csvfile in csvlist:
        os.remove(csvfile)

#name the static slides for different type
def name_slide(type):
    if type == "DESfile_axe_path":
        name = "AxeBlade Static"
    elif type == "DESfile_pdc_path":
        name = "PDC Static"
    elif type == "DESfile_hyper_path":
        name = "HyperBlade Static"
    elif type == "DESfile_stinger_path":
        name = "StingBlade Static"
    elif type == "DESfile_E616_path":
        name = '"EverythingBlade" Static'
    elif type == "DESfile_echo_path":
        name = "StrataBlade Static"
    elif type == "DESfile_px_path":
        name = "PX Static"
    else:
        name = type.split("_")[1]
    return name

def gengerate_profile(parent_dir, template_path, config_path):
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    prs = Presentation(template_path)
#find the "Radius" & "Height" column
    try:
        row_title1 = pd.read_csv(parent_dir + 'BIT1.csv', header=1)
        row_title2 = pd.read_csv(parent_dir + "BIT2.csv", header=1)
        # pd.read_csv('BIT1.csv', delim_whitespace=True, index_col="name")
        # pd.read_csv(parent_dir + "BIT2.csv" ,usecols=[0])
        # row_title = pd.read_csv(parent_dir + "BIT1.csv", nrows = 5)
        x_axis = row_title1[["Radius"]]
        y1_axis = row_title1[["Height"]]
        y2_axis = row_title2[["Height"]]
    except KeyError:
        row_title1 = pd.read_csv(parent_dir + 'BIT1.csv', header=3)
        row_title2 = pd.read_csv(parent_dir + "BIT2.csv", header=3)
        # pd.read_csv('BIT1.csv', delim_whitespace=True, index_col="name")
        # pd.read_csv(parent_dir + "BIT2.csv" ,usecols=[0])
        # row_title = pd.read_csv(parent_dir + "BIT1.csv", nrows = 5)
        x_axis = row_title1[["Radius"]]
        y1_axis = row_title1[["Height"]]
        y2_axis = row_title2[["Height"]]
    else:
        logging.warning("Invalid BIT file, pls check.")
        
    pic_left = int(prs.slide_width*0.13)
    pic_top = int(prs.slide_height*0.1)
    pic_width = int(prs.slide_width*0.65)
    pic_height = int(prs.slide_height*0.8)
    time_title1 = cf.get("profile_title", "title1")
    time_title2 = cf.get("profile_title", "title2")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Profile Comparison"
    x1 = x_axis
    y1 = y1_axis
    #close axis
    plt.axis('off')
    # plotting the line1's points 
    plt.plot(x1, y1, label = "build" + time_title1, color='r', linewidth=2)
    # line2's points
    x2 = x_axis
    y2 = y2_axis
    # plotting the line2's points 
    plt.plot(x2, y2, label = "build" + time_title2, color='b', linewidth=2)
    # naming the x&y axis
    x_axis_label = "Radius(in)"
    y_axis_label = "Height(in)"
    plot_title = "Profile Comparison"
    plt.xlabel(x_axis_label)
    plt.ylabel(y_axis_label)
    # giving a title to my graph
    plt.title(plot_title, loc='center')
    plt.legend()
    buffer_png = io.BytesIO()
    plt.savefig(buffer_png, format = "png")
    pic1 = slide.shapes.add_picture(buffer_png, pic_left, pic_top, pic_width, pic_height)
    plt.close()
    buffer_png.close()
    prs.save(template_path)

def bkrake_sdrake(parent_dir, template_path, config_path):
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    prs = Presentation(template_path)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Back Rake/Side Rake Information"
    for side_type in ['Back_Rake','Side_Rake']:
    #find the "Radius" & "Height" column
        pic_left_list = [int(prs.slide_width*0.1), int(prs.slide_width*0.5)]
        try:
            row_title1 = pd.read_csv(parent_dir + 'BIT1.csv', header=1)
            row_title2 = pd.read_csv(parent_dir + "BIT2.csv", header=1)
            # pd.read_csv('BIT1.csv', delim_whitespace=True, index_col="name")
            # pd.read_csv(parent_dir + "BIT2.csv" ,usecols=[0])
            # row_title = pd.read_csv(parent_dir + "BIT1.csv", nrows = 5)
            x_axis = row_title1[["Radius"]]
        except KeyError:
            row_title1 = pd.read_csv(parent_dir + 'BIT1.csv', header=3)
            row_title2 = pd.read_csv(parent_dir + "BIT2.csv", header=3)
            # pd.read_csv('BIT1.csv', delim_whitespace=True, index_col="name")
            # pd.read_csv(parent_dir + "BIT2.csv" ,usecols=[0])
            # row_title = pd.read_csv(parent_dir + "BIT1.csv", nrows = 5)
            x_axis = row_title1[["Radius"]]
        else:
            logging.warning("Invalid BIT file, pls check.")
        if side_type == "Back_Rake":
            x_axis = row_title1[["Radius"]]
            y1_axis = row_title1[["Back_Rake"]]
            y2_axis = row_title2[["Back_Rake"]]
            pic_left = pic_left_list[0]
            # naming the x&y axis
            x_axis_label = "Radius(in)"
            y_axis_label = "Back Rake Angle (deg)"
            plot_title = "Back Rake"
        elif side_type == "Side_Rake":
            x_axis = row_title1[["Radius"]]
            y1_axis = row_title1[["Side_Rake"]]
            y2_axis = row_title2[["Side_Rake"]]
            pic_left = pic_left_list[1]
            # naming the x&y axis
            x_axis_label = "Radius(in)"
            y_axis_label = "Side Rake Angle (deg)"
            plot_title = "Side Rake"
        pic_top = int(prs.slide_height*0.1)
        pic_width = int(prs.slide_width*0.4)
        pic_height = int(prs.slide_height*0.7)
        time_title1 = cf.get("profile_title", "title1")
        time_title2 = cf.get("profile_title", "title2")
        x1 = x_axis
        y1 = y1_axis
        plt.xticks([])
        # plotting the line1's points 
        plt.plot(x1, y1, label = "build" + time_title1, color='r', linewidth=2)
        # line2's points
        x2 = x_axis
        y2 = y2_axis
        # plotting the line2's points 
        plt.plot(x2, y2, label = "build" + time_title2, color='b', linewidth=2)
        plt.xlabel(x_axis_label)
        plt.ylabel(y_axis_label)
        # giving a title to my graph
        plt.title(plot_title, loc='center')
        plt.legend()
        buffer_png = io.BytesIO()
        plt.savefig(buffer_png, format = "png")
        pic1 = slide.shapes.add_picture(buffer_png, pic_left, pic_top, pic_width, pic_height)
        plt.close()
        buffer_png.close()
    prs.save(template_path)

def isStinger(str_cutter_type) :   
    is_stinger = False 
    for type in ["565_790_85_90", "565_690_85_90", "566_790_85_90", "566_690_85_90", "440_540_85_90", "440_537_85_90",
                  "625_826_96_94", "562_7880_S", "B565_690_441_900_80", "B565_790_541_900_80", "New_8590_L", "562_9694_L",
                   "562_8590_S", "562_7880_L", "562_8590_L",]:
        if str_cutter_type == type:
            is_stinger = True
            break
    return is_stinger

def cutter_info(parent_dir, template_path, config_path):
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    prs = Presentation(template_path)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Cutter Information"
    #find the "Radius" & "Height" column
    try:
        row_title1 = pd.read_csv(parent_dir + 'BIT1.csv', header=1)
        row_title2 = pd.read_csv(parent_dir + "BIT2.csv", header=1)
        # pd.read_csv('BIT1.csv', delim_whitespace=True, index_col="name")
        # pd.read_csv(parent_dir + "BIT2.csv" ,usecols=[0])
        # row_title = pd.read_csv(parent_dir + "BIT1.csv", nrows = 5)
        x_axis = row_title1[["Radius"]]
    except KeyError:
        row_title1 = pd.read_csv(parent_dir + 'BIT1.csv', header=3)
        row_title2 = pd.read_csv(parent_dir + "BIT2.csv", header=3)
        # pd.read_csv('BIT1.csv', delim_whitespace=True, index_col="name")
        # pd.read_csv(parent_dir + "BIT2.csv" ,usecols=[0])
        # row_title = pd.read_csv(parent_dir + "BIT1.csv", nrows = 5)
        x_axis = row_title1[["Radius"]]
    else:
        logging.warning("Invalid BIT file, pls check.")
    # pd.read_csv('BIT1.csv', delim_whitespace=True, index_col="name")
    # pd.read_csv(parent_dir + "BIT2.csv" ,usecols=[0])
    # row_title = pd.read_csv(parent_dir + "BIT1.csv", nrows = 5)
    x_axis = row_title1[["Radius"]]
    y1_axis = row_title1[["Cutter_Distance"]]
    y2_axis = row_title2[["Cutter_Distance"]]
    # naming the x&y axis
    x_axis_label = "Radius (in)"
    y_axis_label = "Distance (in)"
    plot_title = "Cutter Distance along Profile"
    pic_left_list = [int(prs.slide_width*0.13), int(prs.slide_width*0.55)]
    pic_top = int(prs.slide_height*0.1)
    pic_width = int(prs.slide_width*0.4)
    pic_height = int(prs.slide_height*0.5)
    time_title1 = cf.get("profile_title", "title1")
    time_title2 = cf.get("profile_title", "title2")
    x1 = x_axis
    y1 = y1_axis
    plt.grid()
    # plotting the line1's points 
    plt.plot(x1, y1, label = "build" + time_title1, color='r', linewidth=2)
    # line2's points
    x2 = x_axis
    y2 = y2_axis
    # plotting the line2's points 
    plt.plot(x2, y2, label = "build" + time_title2, color='b', linewidth=2)
    plt.xlabel(x_axis_label)
    plt.ylabel(y_axis_label)
    # giving a title to my graph
    plt.title(plot_title, loc='center')
    plt.legend()
    buffer_png = io.BytesIO()
    plt.savefig(buffer_png, format = "png")
    pic1 = slide.shapes.add_picture(buffer_png, pic_left_list[0], pic_top, pic_width, pic_height)
    plt.close()
    buffer_png.close()
    
    # the cutter count bar
    # row_title = pd.read_csv(parent_dir + 'INFO.xlsx', header=2)
    xl = win32.Dispatch("Excel.Application")
    xl.DisplayAlerts = False
    # xl.Visible = None
    xlbook = xl.Workbooks.Open(Filename = parent_dir + "INFO.xlsx", ReadOnly = True)
    # xl.Active.Visible = False
    sheet = xlbook.Worksheets("INFO")
    # is stinger or not
    cutter_type_list = []
    stinger = False
    for i in range(4):
        cutter_type_list.append(sheet.Range("I" + str(i+4)).Value.split(",")[0].split("(")[0].strip())
    for cutter_type in cutter_type_list:
        if isStinger(cutter_type) :
            stinger = True
            break
    # bit name
    bit_name = sheet.Range("A4").Value
    cone_num1 = []
    cone_num2 = []
    cone_num3 = []
    cone_num4 = []
    pic_top = int(prs.slide_height*0.1)
    pic_width = int(prs.slide_width*0.4)
    pic_height = int(prs.slide_height*0.5)
    i = 0
    # xl.ActiveSheet.Cells(4, iCol).Value
    while i <= 4:
        cone_num1.append(sheet.Range("L" + str(i+4)).Value)
        cone_num2.append(sheet.Range("L" + str(i+5)).Value)
        cone_num3.append(sheet.Range("L" + str(i+6)).Value)
        cone_num4.append(sheet.Range("L" + str(i+7)).Value)
        i += 4
    bottom_num1 = []
    bottom_num2 = []
    for i in range(0, len(cone_num1)):
        bottom_num1.append(cone_num1[i] + cone_num2[i])
    for i in range(0, len(cone_num1)):
        bottom_num2.append(bottom_num1[i] + cone_num3[i])
    name_list = ["build" + time_title1, "build" + time_title2]
    N = 2
    ind = np.arange(N)
    plt.grid(axis="y")
    plt.xticks(ind, (name_list[0], name_list[1]))
    plt.bar(ind, cone_num1, label='Cone',fc = 'g')
    plt.bar(ind, cone_num2, bottom=cone_num1, label='Nose',tick_label = name_list,fc = 'gold')
    plt.bar(ind, cone_num3, bottom=bottom_num1, label='Shoulder',tick_label = name_list,fc = 'b')
    plt.bar(ind, cone_num4, bottom=bottom_num2, label='Gage',tick_label = name_list,fc = 'orange')
    plt.legend()
    buffer_png = io.BytesIO()
    plt.savefig(buffer_png, format = "png")
    pic1 = slide.shapes.add_picture(buffer_png, pic_left_list[1], pic_top, pic_width, pic_height)
    plt.close()
    xl.DisplayAlerts = False
    xlbook.Close()
    buffer_png.close()

    # Cutting Structure Info sheet
    # mpl.rcParams['axes.unicode_minus'] = False
    # fig = plt.figure(figsize=(12.7, 3), dpi=900)
    # ax = fig.add_subplot(111, frame_on=False,)
    # # hiden axis
    # ax.xaxis.set_visible(False)  # hide the x axis
    # ax.yaxis.set_visible(False)  # hide the y axis
    # # read excel
    # datas = pd.read_excel(parent_dir + 'INFO.xlsx')
    # datas = datas.iloc[1:10, 0:15]
    # # print(datas)
    # # generate pic
    # table(ax, datas, loc='center')  # where df is your data frame
    # # save
    pic_top = int(prs.slide_height*0.6)
    pic_width = int(prs.slide_width*0.8)
    pic_height = int(prs.slide_height*0.3)
    xl.DisplayAlerts = False
    xlbook = xl.Workbooks.Open(Filename = parent_dir + "INFO.xlsx", ReadOnly = True)
    # xlbook.Visible = False
    sheet = xlbook.Worksheets("INFO")
    excel_range = sheet.Range("A2:P11")
    excel_range.Copy()
    img = ImageGrab.grabclipboard()
    # img = Image.open(img_path)
    with BytesIO() as out:
        img.save(out, format='JPEG')
        # with open('out_stream.jpg', 'wb') as fo:
        #     fo.write(out.getvalue())
    # buffer_png = io.BytesIO()
    # img.save(buffer_png)
    # img.savefig(img, format = "jpg")
        pic1 = slide.shapes.add_picture(out, pic_left_list[0], pic_top, pic_width, pic_height)
    out.close()
    xl.DisplayAlerts = False
    xlbook.Close()
    # img.close()
    # xl.Visible = True
    
    #Cutting Structure Info2 sheet
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Cutter Information"
    pic_top = int(prs.slide_height*0.1)
    pic_width = int(prs.slide_width*0.4)
    pic_height = int(prs.slide_height*0.8)
    xl.DisplayAlerts = False
    xlbook = xl.Workbooks.Open(Filename = parent_dir + "INFO2.xlsx", ReadOnly = True)
    # xl.Visible = False
    sheet = xlbook.Worksheets("INFO2")
    excel_range = sheet.Range("A1:D21")
    excel_range.Copy()
    img = ImageGrab.grabclipboard()
    # img = Image.open(img_path)
    with BytesIO() as out:
        img.save(out, format='JPEG')
        # with open('out_stream.jpg', 'wb') as fo:
        #     fo.write(out.getvalue())
    # buffer_png = io.BytesIO()
    # img.save(buffer_png)
    # img.savefig(img, format = "jpg")
        pic1 = slide.shapes.add_picture(out, pic_left_list[0], pic_top, pic_width, pic_height)
    out.close()
    xl.DisplayAlerts = False
    xlbook.Close()
    xl.DisplayAlerts = True
    prs.save(template_path)
    return stinger, bit_name

                
def static_result_plots(parent_dir, template_path, config_path, bit_info):
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    prs = Presentation(template_path)
    
    
    with open(parent_dir + 'SUMMARY.csv', encoding='utf-8') as case_file:
        reader = csv.reader(case_file)
        case_list = [row[1] for row in reader]
    
    case_list = list(filter(None, case_list))
    case_list.remove(case_list[0])
    #print(case_list)
    # case_file = pd.read_csv(parent_dir + 'SUMMARY.csv', header=1, skip_blank_lines=True)
    # # case_file = pd.read_csv(parent_dir + 'SUMMARY.csv', header=2)
    # for rows in case_file:
    #     print(rows)
    # case_list = case_file[bit_info[0]][0]
    # case_list = case_file[[bit_info[1]]]
    # stinger
    if bit_info[0]:
        icase = 0
        for case in case_list:
            for perc in ['0.5', '0.75', '0.95']:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Static Analysis Results"
                # caseplot_num = pd.read_csv(parent_dir + "BIT1.csv", nrows = 1)
                i = 0
                for caseplot in ['Circumferential Forces', 'Workrate of Circumferential Forces', 'Normal Forces', 'Workrate of Normal Forces']:
                    #find the "Radius" & "Height" column
                    # with open(parent_dir + 'SUMMARY.csv',encoding='utf-8') as case_file:
                    #     reader = csv.reader(case_file)
                    #     case_list = [row[1] for row in reader]
                    row_title1 = pd.read_csv(parent_dir + 'BIT1.csv', header=1)
                    row_title2 = pd.read_csv(parent_dir + "BIT2.csv", header=1)
                    # pd.read_csv('BIT1.csv', delim_whitespace=True, index_col="name")
                    # pd.read_csv(parent_dir + "BIT2.csv" ,usecols=[0])
                    # row_title = pd.read_csv(parent_dir + "BIT1.csv", nrows = 5)
                    x_axis = row_title1[["Radius"]]
                    # y1_axis = row_title1[[perc]][ icase*4]
                    if i == 0 and icase == 0:
                        y1_axis = row_title1[[perc]]
                        y2_axis = row_title2[[perc]]
                    else:
                        y1_axis = row_title1[[perc+'.'+ str(i+icase*4)]]
                        y2_axis = row_title2[[perc+'.'+ str(i+icase*4)]]
                    # naming the x&y axis
                    x_axis_label = "Radius (in)"
                    if caseplot == 'Circumferential Forces':
                        y_axis_label = "Circumferential Forces (lbf) "
                        pic_left_list = int(prs.slide_width*0.5)
                        pic_top = int(prs.slide_height*0.1)
                    elif caseplot == 'Workrate of Circumferential Forces':
                        y_axis_label = "Workrate (lbf-ft/s)"
                        pic_left_list = int(prs.slide_width*0.5)
                        pic_top = int(prs.slide_height*0.5)
                    elif caseplot == 'Normal Forces':
                        y_axis_label = "Normal Forces (lbf)"
                        pic_left_list = int(prs.slide_width*0.2)
                        pic_top = int(prs.slide_height*0.1)
                    elif caseplot == 'Workrate of Normal Forces':
                        y_axis_label = "Workrate (lbf-ft/s)"
                        pic_left_list = int(prs.slide_width*0.2)
                        pic_top = int(prs.slide_height*0.5)
                    plot_title = caseplot + "   " + perc + "\n" + case
                    pic_width = int(prs.slide_width*0.3)
                    pic_height = int(prs.slide_height*0.4)
                    time_title1 = cf.get("profile_title", "title1")
                    time_title2 = cf.get("profile_title", "title2")
                    x1 = x_axis
                    y1 = y1_axis
                    plt.grid(axis="y")
                    # plotting the line1's points 
                    plt.plot(x1, y1, label = "build" + time_title1, color='r', linewidth=2)
                    # line2's points
                    x2 = x_axis
                    y2 = y2_axis
                    # plotting the line2's points 
                    plt.plot(x2, y2, label = "build" + time_title2, color='b', linewidth=2)
                    plt.xlabel(x_axis_label)
                    plt.ylabel(y_axis_label)
                    # giving a title to my graph
                    plt.title(plot_title, loc='center')
                    plt.legend()
                    buffer_png = io.BytesIO()
                    plt.savefig(buffer_png, format = "png")
                    pic1 = slide.shapes.add_picture(buffer_png, pic_left_list, pic_top, pic_width, pic_height)
                    plt.close()
                    buffer_png.close()
                    i += 1
            icase += 1
        prs.save(template_path)
        #add comparsion slides:
        i = 0
        comparison_list = ['Instantaneous ROP (ft/hr)', 'Torque (ft-lbf)', 'Bit Advance (in/rev)', r'Total imbalance force (% of WOB)', 'Beta Angle', 
                         r'Circumferential imbalance force (% of WOB)', r'Radial imbalance force (% of WOB)']
        for caseplot in ['ROP Comparison', 'Torque Comparison', 'Bit Advance Comparison', 'Total Imbalance Force Comparison', 'Beta Angle Comparison', 
                         'Circumferential Imbalance Force Comparison', 'Radial Imbalance Force Comparison']:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = caseplot
            if caseplot == 'ROP Comparison' or caseplot == 'Torque Comparison' or caseplot == 'ROP/Torque Comparison':
                pic_left = int(prs.slide_width*0.05)
                pic_top = int(prs.slide_height*0.3)
                pic_width = int(prs.slide_width*0.45)
                pic_height = int(prs.slide_height*0.4)
                # for compr_case in case_list:
                icase = 0
                with open(parent_dir + 'SUMMARY.csv', encoding='utf-8') as case_file:
                    reader = csv.reader(case_file)
                    for row in reader:
                        if icase == 1 :
                            case_column = row.index(comparison_list[i])
                            if caseplot =='ROP Comparison':
                                case_change = row.index(r'% change in Instantaneous ROP')
                            elif caseplot =='Torque Comparison':
                                case_change = row.index(r'% change in Instantaneous Torque')
                            elif caseplot =='ROP/Torque Comparison':
                                case_change = row.index(r'% Change in ROP/Torque')
                            break
                        icase += 1
                    row_title1 = [row[case_column + 1] for row in reader]
                    row_title1 = list(filter(None, row_title1))
                    row_title1 = row_title1[1:]
                    case_file.seek(0)
                    reader = csv.reader(case_file)
                    row_title2 = [row[case_column + 2] for row in reader]
                    row_title2 = list(filter(None, row_title2))
                    row_title2 = row_title2[1:]
                    case_file.seek(0)
                    reader = csv.reader(case_file)
                    row_title3 = [row[case_change + 1] for row in reader]
                    row_title3 = list(filter(None, row_title3))
                    row_title3 = row_title3[1:]
                y1 = []
                y2 = []
                y3 = []
                for item in row_title1:
                    y1.append(float(item))
                for item in row_title2:
                    y2.append(float(item))
                for item in row_title3:
                    y3.append(float(item))  
                # y3 = [0] * len(y1)
                # for inum in range(len(y1)):
                #     y3[inum] = (y2[inum]-y1[inum])/y1[inum]
                # x = np.array(case_list)
                x = np.arange(len(case_list))
                time_title1 = cf.get("profile_title", "title1")
                time_title2 = cf.get("profile_title", "title2")

                # y1 = np.array([3, 8, 1, 10])
                plt.grid(axis="y")
                plt.bar(x, y1, color='r', width = 0.4, label = "build" + time_title1)
                plt.bar(x+0.4, y1, color='b', width = 0.4, label = "build" + time_title2)
                plt.xticks(x + 0.4, case_list)
                # plt.xlabel(x_axis_label)
                # plt.ylabel(y_axis_label)
                plt.title(comparison_list[i], loc='center')
                plt.legend()
                buffer_png = io.BytesIO()
                plt.savefig(buffer_png, format = "png")
                pic1 = slide.shapes.add_picture(buffer_png, pic_left, pic_top, pic_width, pic_height)
                plt.close()
                buffer_png.close()

                #add another plot in the three caseplots:
                
                x = np.array(case_list)
                # y1 = np.array([3, 8, 1, 10])
                plt.grid(axis="y")
                plt.bar(x, y3, color='r', width = 0.4, label = "build" + time_title2)
                # plt.xlabel(x_axis_label)
                # plt.ylabel(y_axis_label)

                plt.title(plot_title, loc='center')
                plt.legend()
                buffer_png = io.BytesIO()
                plt.savefig(buffer_png, format = "png")
                pic1 = slide.shapes.add_picture(buffer_png, pic_left, pic_top, pic_width, pic_height)
                plt.close()
                buffer_png.close()
            else:
                pic_left = int(prs.slide_width*0.1)
                pic_top = int(prs.slide_height*0.1)
                pic_width = int(prs.slide_width*0.8)
                pic_height = int(prs.slide_height*0.8)
                # for compr_case in case_list:
                icase = 0
                with open(parent_dir + 'SUMMARY.csv', encoding='utf-8') as case_file:
                    reader = csv.reader(case_file)
                    for row in reader:
                        if icase == 1 :
                            case_column = row.index(comparison_list[i])
                            break
                        icase += 1
                    row_title1 = [row[case_column + 1] for row in reader]
                    row_title1 = list(filter(None, row_title1))
                    row_title1 = row_title1[1:]
                    y1 = []
                    for item in row_title1:
                        y1.append(float(item))
                    case_file.seek(0)
                    reader = csv.reader(case_file)
                    row_title2 = [row[case_column + 2] for row in reader]
                    row_title2 = list(filter(None, row_title2))
                    row_title2 = row_title2[1:]
                    y2 = []
                    for item in row_title2:
                        y2.append(float(item))
                    # y3 = [0] * len(y1)
                    # for inum in range(len(y1)):
                    #     y3[inum] = (y2[inum]-y1[inum])/y1[inum]
                    x = np.arange(len(case_list))
                    time_title1 = cf.get("profile_title", "title1")
                    time_title2 = cf.get("profile_title", "title2")
                    mpl.rcParams['font.size'] = 10
                    # y1 = np.array([3, 8, 1, 10])
                    plt.grid(axis="y")
                    plt.bar(x, y1, color='r', width = 0.4, label = "build" + time_title1)
                    plt.bar(x+0.4, y1, color='b', width = 0.4, label = "build" + time_title2)
                    plt.xticks(x + 0.4, case_list)
                    # plt.xlabel(x_axis_label)
                    # plt.ylabel(y_axis_label)
                    plt.title(comparison_list[i], loc='center')
                    plt.legend()
                    buffer_png = io.BytesIO()
                    plt.savefig(buffer_png, format = "png")
                    pic1 = slide.shapes.add_picture(buffer_png, pic_left, pic_top, pic_width, pic_height)
                    plt.close()
                    buffer_png.close()
            i += 1
    # not stinger
    else:
        icase = 0
        for case in case_list:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Static Analysis Results"
            # caseplot_num = pd.read_csv(parent_dir + "BIT1.csv", nrows = 1)
            i = 0
            for caseplot in ['Cut Area', 'Cut Volume', 'Circumferential Forces', 'Normal Forces', 'Workrate of Circumferential Forces', 'Workrate of Normal Forces']:
                #find the "Radius" & "Height" column
                icase = 0
                with open(parent_dir + 'BIT1.csv', encoding='utf-8') as case_file:
                    reader = csv.reader(case_file)
                    for row in reader:
            #save rop as a list
                        if icase == 1 :
                            # first_line = str(row).split(',')
                            case_column = row.index(case)
                            break
                        icase += 1
                    if i == 0 or i == 1:
                        row_title1 = [row[case_column + i] for row in reader]
                    elif i == 2:
                        row_title1 = [row[case_column + 3] for row in reader]
                    elif i == 3:
                        row_title1 = [row[case_column + 5] for row in reader]
                    elif i == 4:
                        row_title1 = [row[case_column + 15] for row in reader]
                    elif i == 5:
                        row_title1 = [row[case_column + 16] for row in reader]
                    #remove the not required element, under the row[1]
                    row_title1 = row_title1[2:]
                    #remove '' in list 
                    row_title1 = list(filter(None, row_title1))
                    case_file.seek(0)
                    reader = csv.reader(case_file)
                    x_axis = [row[1] for row in reader]
                    x_axis = list(filter(None, x_axis))
                    x_axis = x_axis[1:]
                icase = 0
                with open(parent_dir + 'BIT2.csv', encoding='utf-8') as case_file:
                    reader = csv.reader(case_file)
                    for row in reader:
            #save rop as a list
                        if icase == 1 :
                            # first_line = str(row).split(',')
                            # case_column = first_line.index(case)
                            case_column = row.index(case)
                            break
                        icase += 1
                    if i == 0 or i == 1:
                        row_title2 = [row[case_column + i] for row in reader]
                    elif i == 2:
                        row_title2 = [row[case_column + 3] for row in reader]
                    elif i == 3:
                        row_title2 = [row[case_column + 5] for row in reader]
                    elif i == 4:
                        row_title2 = [row[case_column + 15] for row in reader]
                    elif i == 5:
                        row_title2 = [row[case_column + 16] for row in reader]
                    #remove the not required element
                    row_title2 = row_title2[2:]
                    #remove '' in list 
                    row_title2 = list(filter(None, row_title2))
                y1_axis = []
                y2_axis = []
                for item in row_title1:
                    y1_axis.append(float(item))
                for item in row_title2:
                    y2_axis.append(float(item))
                # naming the x&y axis
                x_axis_label = "Radius (in)"
                if caseplot == 'Circumferential Forces':
                    y_axis_label = "Circumferential Forces (lbf) "
                    pic_left = int(prs.slide_width*0.37)
                    pic_top = int(prs.slide_height*0.1)
                elif caseplot == 'Workrate of Circumferential Forces':
                    y_axis_label = "Workrate (lbf-ft/s)"
                    pic_left = int(prs.slide_width*0.37)
                    pic_top = int(prs.slide_height*0.5)
                elif caseplot == 'Normal Forces':
                    y_axis_label = "Normal Forces (lbf)"
                    pic_left = int(prs.slide_width*0.12)
                    pic_top = int(prs.slide_height*0.1)
                elif caseplot == 'Workrate of Normal Forces':
                    y_axis_label = "Workrate (lbf-ft/s)"
                    pic_left = int(prs.slide_width*0.12)
                    pic_top = int(prs.slide_height*0.5)
                elif caseplot == 'Cut Area':
                    y_axis_label = 'Cut Area($\mathregular{in^3}$)'
                    pic_left = int(prs.slide_width*0.62)
                    pic_top = int(prs.slide_height*0.1)
                elif caseplot == 'Cut Volume':
                    y_axis_label = 'Cut Volume($\mathregular{in^3}$)'
                    pic_left = int(prs.slide_width*0.62)
                    pic_top = int(prs.slide_height*0.5)
                plot_title = caseplot + "\n" + case
                pic_width = int(prs.slide_width*0.25)
                pic_height = int(prs.slide_height*0.4)
                time_title1 = cf.get("profile_title", "title1")
                time_title2 = cf.get("profile_title", "title2")
                x_axis = x_axis[:len(y1_axis)]
                x_axis_num = []
                for item in x_axis:
                    x_axis_num.append(float(item))
                x1 = x_axis_num
                y1 = y1_axis
                plt.grid(axis="y")
                # plotting the line1's points 
                plt.plot(x1, y1, label = "build" + time_title1, color='r', linewidth=2)
                # line2's points
                x2 = x_axis_num
                y2 = y2_axis
                # plotting the line2's points 
                plt.plot(x2, y2, label = "build" + time_title2, color='b', linewidth=2)
                plt.xlabel(x_axis_label)
                plt.ylabel(y_axis_label)
                # giving a title to my graph
                plt.title(plot_title, loc='center')
                plt.legend()
                buffer_png = io.BytesIO()
                plt.savefig(buffer_png, format = "png")
                pic1 = slide.shapes.add_picture(buffer_png, pic_left, pic_top, pic_width, pic_height)
                plt.close()
                buffer_png.close()
                i += 1
            icase += 1
        prs.save(template_path)
        #add comparsion slides:
        i = 0
        comparison_list = ['Instantaneous ROP (ft/hr)', 'Torque (ft-lbf)', 'ROP/Torque', 'Bit Advance (in/rev)', r'Total imbalance force (% of WOB)', 'Beta Angle', 
                         r'Circumferential imbalance force (% of WOB)', r'Radial imbalance force (% of WOB)', 'Normalized IF Ratio']
        for caseplot in ['ROP Comparison', 'Torque Comparison', 'ROP/Torque Comparison', 'Bit Advance Comparison', 'Total Imbalance Force Comparison', 'Beta Angle Comparison', 
                         'Circumferential Imbalance Force Comparison', 'Radial Imbalance Force Comparison', 'Normalized Imbalance Force Ratio Comparison']:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = caseplot
            if caseplot == 'ROP Comparison' or caseplot == 'Torque Comparison' or caseplot == 'ROP/Torque Comparison':
                pic_left = int(prs.slide_width*0.05)
                pic_top = int(prs.slide_height*0.3)
                pic_width = int(prs.slide_width*0.45)
                pic_height = int(prs.slide_height*0.4)
                # for compr_case in case_list:
                icase = 0
                with open(parent_dir + 'SUMMARY.csv', encoding='utf-8') as case_file:
                    reader = csv.reader(case_file)
                    for row in reader:
                        if icase == 1 :
                            case_column = row.index(comparison_list[i])
                            if caseplot =='ROP Comparison':
                                case_change = row.index(r'% change in Instantaneous ROP')
                            elif caseplot =='Torque Comparison':
                                case_change = row.index(r'% change in Instantaneous Torque')
                            elif caseplot =='ROP/Torque Comparison':
                                case_change = row.index(r'% Change in ROP/Torque')
                            break
                        icase += 1
                    row_title1 = [row[case_column + 1] for row in reader]
                    row_title1 = list(filter(None, row_title1))
                    row_title1 = row_title1[1:]
                    case_file.seek(0)
                    reader = csv.reader(case_file)
                    row_title2 = [row[case_column + 2] for row in reader]
                    row_title2 = list(filter(None, row_title2))
                    row_title2 = row_title2[1:]
                    case_file.seek(0)
                    reader = csv.reader(case_file)
                    row_title3 = [row[case_change + 1] for row in reader]
                    row_title3 = list(filter(None, row_title3))
                    row_title3 = row_title3[1:]
                y1 = []
                y2 = []
                y3 = []
                for item in row_title1:
                    y1.append(float(item))
                for item in row_title2:
                    y2.append(float(item))
                for item in row_title3:
                    y3.append(float(item))  
                # y3 = [0] * len(y1)
                # for inum in range(len(y1)):
                #     y3[inum] = (y2[inum]-y1[inum])/y1[inum]
                # x = np.array(case_list)
                x = np.arange(len(case_list))
                time_title1 = cf.get("profile_title", "title1")
                time_title2 = cf.get("profile_title", "title2")
                mpl.rcParams['font.size'] = 10
                # y1 = np.array([3, 8, 1, 10])
                plt.grid(axis="y")
                plt.bar(x, y1, color='r', width = 0.4, label = "build" + time_title1)
                plt.bar(x+0.4, y1, color='b', width = 0.4, label = "build" + time_title2)
                plt.xticks(x + 0.4, case_list)
                # plt.xlabel(x_axis_label)
                # plt.ylabel(y_axis_label)
                plt.title(comparison_list[i], loc='center')
                plt.legend()
                buffer_png = io.BytesIO()
                plt.savefig(buffer_png, format = "png")
                pic1 = slide.shapes.add_picture(buffer_png, pic_left, pic_top, pic_width, pic_height)
                plt.close()
                buffer_png.close()
                prs.save(template_path)
                
                #add another plot in the three caseplots:
                pic_left = int(prs.slide_width*0.5)
                x = np.arange(len(case_list))
                # y1 = np.array([3, 8, 1, 10])
                plt.grid(axis="y")
                plt.bar(x, y3, color='b', width = 0.4, label = "build" + time_title2)
                plt.xticks(x + 0.4, case_list)
                # plt.xlabel(x_axis_label)
                # plt.ylabel(y_axis_label)

                plt.title(plot_title, loc='center')
                plt.legend()
                buffer_png = io.BytesIO()
                plt.savefig(buffer_png, format = "png")
                pic1 = slide.shapes.add_picture(buffer_png, pic_left, pic_top, pic_width, pic_height)
                plt.close()
                buffer_png.close()
            else:
                pic_left = int(prs.slide_width*0.1)
                pic_top = int(prs.slide_height*0.1)
                pic_width = int(prs.slide_width*0.8)
                pic_height = int(prs.slide_height*0.8)
                # for compr_case in case_list:
                icase = 0
                with open(parent_dir + 'SUMMARY.csv', encoding='utf-8') as case_file:
                    reader = csv.reader(case_file)
                    for row in reader:
                        if icase == 1 :
                            case_column = row.index(comparison_list[i])
                            break
                        icase += 1
                    row_title1 = [row[case_column + 1] for row in reader]
                    row_title1 = list(filter(None, row_title1))
                    row_title1 = row_title1[1:]
                    y1 = []
                    for item in row_title1:
                        y1.append(float(item))
                    case_file.seek(0)
                    reader = csv.reader(case_file)
                    row_title2 = [row[case_column + 2] for row in reader]
                    row_title2 = list(filter(None, row_title2))
                    row_title2 = row_title2[1:]
                    y2 = []
                    for item in row_title2:
                        y2.append(float(item))
                    # y3 = [0] * len(y1)
                    # for inum in range(len(y1)):
                    #     y3[inum] = (y2[inum]-y1[inum])/y1[inum]
                    # x = np.array(case_list)
                    time_title1 = cf.get("profile_title", "title1")
                    time_title2 = cf.get("profile_title", "title2")
                    # y1 = np.array([3, 8, 1, 10])
                    plt.grid(axis="y")
                    mpl.rcParams['font.size'] = 10
                    x = np.arange(len(case_list))
                    # plt.bar(x, y1, color='r', width = 0.4, label = "build" + time_title1)
                    # plt.bar(x, y1, color='b', width = 0.4, label = "build" + time_title2)
                    plt.bar(x, y1, color='r', width = 0.4, label = "build" + time_title1)
                    plt.bar(x+0.4, y1, color='b', width = 0.4, label = "build" + time_title2)
                    plt.xticks(x + 0.4, case_list)
                    # plt.xlabel(x_axis_label)
                    # plt.ylabel(y_axis_label)
                    plt.title(comparison_list[i], loc='center')
                    plt.legend()
                    buffer_png = io.BytesIO()
                    plt.savefig(buffer_png, format = "png")
                    pic1 = slide.shapes.add_picture(buffer_png, pic_left, pic_top, pic_width, pic_height)
                    plt.close()
                    buffer_png.close()
            i += 1
    prs.save(template_path)

#read bit by vba
def readbit(config_path, static_type, parent_dir):
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    xl = win32.Dispatch("Excel.Application")
    xl.DisplayAlerts = False
    xl.Visible = False
    parent_dir = parent_dir.replace(parent_dir.split("\\")[-1], '')
    xlbook = xl.Workbooks.Open(Filename = parent_dir + "test_automation\\SAM_v11.78.xlsm", ReadOnly = True)
    WorkbookName = "SAM_v11.78.xlsm"
    xl.Workbooks("SAM_v11.78.xlsm").Activate
    MacroName1 = "Module1.ReadBit"
    MacroName3 = "Module1.Clear_All"

    argument2 = False
    argument3 = 1
    argument4 = True
    argument6 = 2
    path_array = []
    # cases_type = ReadIniFileString("static_case_type", "type_list")
    path_array.append(cf.get(static_type, "cases_path1"))
    path_array.append(cf.get(static_type, "cases_path2"))
    xl.Application.Run("'" & WorkbookName & "'!" & MacroName3)
    xl.Application.Run("'" & WorkbookName & "'!" & MacroName1, path_array(0), argument4, argument3, argument2)
    xl.Application.Run("'" & WorkbookName & "'!" & MacroName1, path_array(1), argument4, argument6, argument2)
    xlbook.Close()
    xl.DisplayAlerts = True
    xl.Visible = True

#generate static plots
def generate_static_plots(config_path):
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    template_path = cf.get("merge_template", "dirct")
    prs = Presentation(template_path)
    # clear_csv(config_path)
    static_type_list = str(cf.get("static_case_type", "type_list")).split()
    parent_dir = config_path.replace(config_path.split("\\")[-1], '')
    for static_type in static_type_list:
        readbit(config_path, static_type, parent_dir)
        slide = prs.slides.add_slide(prs.slide_layouts[10])
        slide.shapes.title.text = name_slide(static_type)
        gengerate_profile(parent_dir, template_path, config_path)
        bkrake_sdrake(parent_dir, template_path, config_path)
        bit_info = cutter_info(parent_dir, template_path, config_path)
        static_result_plots(parent_dir, template_path, config_path, bit_info)

    prs.save(template_path)
    
if __name__ == "__main__":
    path = os.getcwd()
    config_ini_dir = path + r'\docs\auto_test_config.ini'
    generate_static_plots(config_ini_dir)