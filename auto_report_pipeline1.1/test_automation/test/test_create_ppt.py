from pptx import Presentation
from pptx.util import Pt
import pandas as pd


class WritePowerPoint:
    def __init__(self, ppt_name, input_excel, title_cover, subtitle):
        self.ppt_name = ppt_name
        self.input_excel = input_excel
        self.title_cover = self.title_per_page = title_cover
        self.subtitle_cover = subtitle
        # 创建空白演示文稿
        self.prs = Presentation()

    def add_cover(self):
        # 添加封面布局幻灯片
        slide_layout_cover = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout_cover)
        # 设置标题和副标题
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = self.title_cover
        subtitle.text = self.subtitle_cover

    def add_slide(self, line2_texts):
        # 添加布局5幻灯片
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        # 设置标题：内容、位置、字体、大小等格式
        title_shape = shapes.title
        title_shape.text = self.title_per_page
        title_shape.left, title_shape.top = Pt(32), Pt(22)
        title_shape.width, title_shape.height = Pt(660), Pt(50)
        tf0 = title_shape.text_frame
        p0 = tf0.paragraphs[0]
        p0.font.name = '微软雅黑'
        p0.font.size = Pt(24)
        # 添加文本框
        left, top, width, height = Pt(32), Pt(82), Pt(665), Pt(396)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame

        def add_paragraph_texts(texts):
            print(texts[0])
            for i, text in enumerate(texts[:-1]):       # 最后的网址作为超链接
                p = tf.add_paragraph()
                p.text = text
                if i == len(texts) - 2:
                    run = p.add_run()
                    run.text = '查看更多'
                    run.hyperlink.address = texts[-1]   # 写入超链接
                    tf.add_paragraph()

        # 在文本框中添加内容
        for lst in line2_texts:
            add_paragraph_texts(lst)

    def run(self):
        # 读取 Excel 数据并进行预处理
        df = pd.read_excel(self.input_excel)
        df['发布机构'] = '发布机构：' + df['发布机构']
        df['发布时间'] = '发布时间：' + df['发布时间']
        df['关键词'] = '关键词：' + df['关键词']

        # 添加封面幻灯片
        self.add_cover()
        # 添加重复格式的幻灯片，每页写 2 条数据
        for i in df.index[::2]:
            self.add_slide([df.loc[i, :].tolist(), df.loc[i+1, :].tolist()])

        # 保存
        self.prs.save(self.ppt_name)


if __name__ == '__main__':
    wpt = WritePowerPoint('news.pptx', 'news.xlsx', 'News Briefs', '2020/4/5')
    wpt.run()