
# From https://www.cnblogs.com/shanger/p/13107188.html
import  os
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# def function_check_sheet(ppt_path):
#     # work_path = r'E:\pyspace\tmp\pptx'
#     # os.chdir(ppt_path)

#     # square_target = [20000, 1000, 1000000, 300]
#     # square_achivement = [19000, 888, 888888, 289]
#     # month_target = [20000, 1000, 100000, 120]
#     # month_achivement = [18000, 888, 888888, 118]
#     # Pass_Fail = []
#     # df = pd.DataFrame(data={'1 季度指标': square_target,
#     #                     '1 季度完成': square_achivement,
#     #                     '1 季度完成率': None,
#     #                     '3 月份指标': month_target,
#     #                     '3 月份完成': month_achivement,
#     #                     '3 月份完成率': None
#     #                     },
#     df = pd.DataFrame(data={'Tool/Program': None,
#                             'Pass/Fail': None,
#                         },
#                     index = ['I-Design (idd) / Scale', 'I-Design (idd) / Modify Cutter/Position',
#                              'IDEAS Static (ids) / Normal Functionality', 'IDEAS Static (ids) / Multi-Core processing',
#                              'IDEAS Dynamic (id3) / Normal Functionality', 'IDEAS Dynamic (id3) / Multi-Core processing',
#                              'IDEAS Autoforce (ida) / Force Balance Optimization', 'IDEAS3 Utility (id3u)',
#                              'IDEAS Analysis Request (iAR) / Create Projects (ideasbuild20200623)', 'Bit Design Toolbox / IAR Project Renamer', 
#                              'Bit Design Toolbox / Rollback IDEAS Project', 'IDM / Create/Run Projectgroup',
#                              'SAM', 'Multi-Bit']
#                     )

#     # # 计算指标完成率，将转换为百分数
#     # df['1 季度完成率'] = (df['1 季度完成'] / df['1 季度指标']).apply(lambda x: format(x, '.1%'))
#     # df['3 月份完成率'] = (df['3 月份完成'] / df['3 月份指标']).apply(lambda x: format(x, '.1%'))


#     # 实例化 ppt 文档对象
#     prs = Presentation(ppt_path)

#     # 插入幻灯片
#     slide = prs.slides.add_slide(prs.slide_layouts[1])
#     slide.shapes.title.text = "Contact Analysis Comparison - Blade Top Contact"

#     # 预设表格总体布局参数
#     rows = 15
#     columns = 2
#     left = Cm(3)   
#     top = Cm(4) 
#     width = Cm(32) 
#     height = Cm(0.1)
#     # 添加表格
#     table = slide.shapes.add_table(rows=rows,
#                             cols=columns,
#                             left=left,
#                             top=top,
#                             width=width,
#                             height=height
#                             )
#     table = table.table

#     # # 调整行高、列宽
#     # for i in range(rows):
#     #     table.rows[i].height = Cm(0.6)
        
#     # for i in range(columns):
#     #     # if i in (1, 5):
#     #     #     continue
#     #     # table.columns[i].width = Cm(3)
#     #     table.columns[i].width = Cm(10)

#     # 写入表头
#     header = df.columns
#     for i, h in enumerate(header):
#         cell = table.cell(0, i)
#         # tf = cell.text_frame
#         # para = tf.add_paragraph()
#         # para.text = h
#         cell.text = h
#         # cell.font.size = Pt(19)
#         # para.alignment = PP_ALIGN.CENTER  # 居中
#         # cell.alignment = PP_ALIGN.LEFT

#     # 写入行名称
#     row_names = df.index
#     for i, r in enumerate(row_names):
#         cell = table.cell(i+1, 0)
#         # tf = cell.text_frame
#         # para = tf.add_paragraph()
#         # para.text = r
#         cell.text = r
#         # para.font.size = Pt(19)
#         # para.alignment = PP_ALIGN.CENTER  # 居中
#         # para.alignment = PP_ALIGN.LEFT
#     prs.save(ppt_path)

def function_check_sheet(ppt_path):
    df = pd.DataFrame(data={'Tool/Program': None,
                            'Pass/Fail': None,
                        },
                    index = ['I-Design (idd) / Scale', 'I-Design (idd) / Modify Cutter/Position',
                             'IDEAS Static (ids) / Normal Functionality', 'IDEAS Static (ids) / Multi-Core processing',
                             'IDEAS Dynamic (id3) / Normal Functionality', 'IDEAS Dynamic (id3) / Multi-Core processing',
                             'IDEAS Autoforce (ida) / Force Balance Optimization', 'IDEAS3 Utility (id3u)',
                             'IDEAS Analysis Request (iAR) / Create Projects (ideasbuild20200623)', 'Bit Design Toolbox / IAR Project Renamer', 
                             'Bit Design Toolbox / Rollback IDEAS Project', 'IDM / Create/Run Projectgroup',
                             'SAM', 'Multi-Bit']
                    )

    prs = Presentation(ppt_path)
    slide = prs.slides[3]
    # slide = prs.slides.add_slide(prs.slide_layouts[1])
    # slide.shapes.title.text = "Contact Analysis Comparison - Blade Top Contact"

    rows = 15
    columns = 2
    left = Cm(3)   
    top = Cm(4) 
    width = Cm(32) 
    height = Cm(0.1)
    table = slide.shapes.add_table(rows=rows,
                            cols=columns,
                            left=left,
                            top=top,
                            width=width,
                            height=height
                            )
    table = table.table
    # colunm_name
    header = df.columns
    for i, h in enumerate(header):
        cell = table.cell(0, i)
        cell.text = h

    # row_name
    row_names = df.index
    for i, r in enumerate(row_names):
        cell = table.cell(i+1, 0)
        cell.text = r
    prs.save(ppt_path)
if __name__ == "__main__":
    function_check_sheet(r"C:\Users\JWang294\src\cis\libcalc\auto_report_pipeline1.1\60095160_RevA_CPC_Template_v4.04 - Copy.pptx")
    