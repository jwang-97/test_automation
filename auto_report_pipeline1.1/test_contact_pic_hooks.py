from pptx import Presentation
import pptx, pptx.util, glob, os, re ,logging
# import scipy.misc
# from pptx.util import Inches
# import imageio
import cv2
import configparser


def contact_pic_hooks(org_import_dir1, org_import_dir2, config_file_path):
    cf = configparser.ConfigParser()
    cf.read(config_file_path, encoding='utf-8')
    template_path = cf.get("merge_template", "dirct")
    prs = Presentation(template_path)
    org_import_dir = (org_import_dir1, org_import_dir2)
    contact_pic_path = ['','']
    contact_pic = ['','']
    i = 0
    while i <= 1:
        # org_path = os.walk(org_import_dir[i+1])
        # for path, dir_list, file_list in org_path:
            # for dir_name in dir_list:
                # if re.search("\Contact\contact_script_test\utility_files", dir_name) != None:
                    # logging.debug('invalid contact path:' + org_path + ", lacks essential utility_files")
        for pic_dir in glob.glob(org_import_dir[i] + r'\Contact\contact_script_test\utility_files*'):
                    # cf.set("Contact_pic_path", "pic_path"+ str(i), os.path.join(path, dir_name))
            contact_pic_path[i] = pic_dir
            break
        i += 1
    pic_left = [int(prs.slide_width*0.08), int(prs.slide_width*0.58)]
    pic_top = int(prs.slide_height*0.3)
    pic_width = int(prs.slide_width*0.4)
    pic_height = int(prs.slide_height*0.6)
    ipath = 0
    time_title1 = cf.get("profile_title", "title1")
    time_title2 = cf.get("profile_title", "title2")
    text_array = ("ROP Controlled: " + "50" + " ft/hr", "Build"+ time_title1, "Build" + time_title2)
    # while ipath <= 1:
    for contact_pic_path1 in glob.glob(contact_pic_path[ipath]+'/*.jpg'):
        if re.search('Blade', contact_pic_path1) == None:
            pic_type = contact_pic_path1.split('\\')[-1]
            bits_num_list = re.findall(r"\d+" ,str(pic_type.split('_')[0]))
            bits_num = bits_num_list[-1]
            # contact_pic[0] = contact_pic_path1
            # pic_path = contact_pic_path1.split('\\')
            # print(pic_path[-1])
            for contact_pic_path2 in glob.glob(contact_pic_path[ipath+1]+'/' + pic_type + '*'):
                # contact_pic[1] = contact_pic_path2
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                # slide.shapes.title.textframe.textrange.text = "test contact slide"
                slide.shapes.title.text = "test contact slide"
                # pic2 = slide.shapes.add_picture()
                img1 = cv2.imread(contact_pic[0])
                img2 = cv2.imread(contact_pic_path[1])
                pic1 = slide.shapes.add_picture(contact_pic_path1, pic_left[0], pic_top, pic_width, pic_height)
                pic2 = slide.shapes.add_picture(contact_pic_path2, pic_left[1], pic_top, pic_width, pic_height)
                print (contact_pic_path1)
                i = 0
                while i <= 2:
                    if i == 0:
                        tb = slide.shapes.add_textbox(300000, 700000, 1, 1)
                    elif i == 1:
                        tb = slide.shapes.add_textbox(1000000, 1600000, 1, 1)
                    else:
                        tb = slide.shapes.add_textbox(8000000, 1600000, 1, 1)
                    p1 =  tb.text_frame.add_paragraph()
                    p1.text = text_array[i]
                    p1.font.size = pptx.util.Pt(25)
                    i+=1

            # img = imageio.imread(g, pilmode="RGB")
            # img = imageio.imread(g, plugin='matplotlib')
            # img = cv2.imread(contact_pic_path1)
            # pic_height = int(pic_width * img.shape[0] / img.shape[1])
            #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
            # pic = slide.shapes.add_picture(contact_pic_path1, pic_left, pic_top, pic_width, pic_height)
            # pic = slide.shapes.add_picture(g, 2.4, 5.87, 16.59, 14.13)
            # pic2 = slide.shapes.add_picture(g, 22.19, 5.87, 16.59, 14.13)

    # prs.save("%s.pptx" % OUTPUT_TAG)
    prs.save(template_path)
if __name__ == "__main__":
    contact_pic_hooks(r"C:\Users\JWang294\Documents\projectfile\test_config_cases\test-20210416173550-lbo-0416cpc", r"C:\Users\JWang294\Documents\projectfile\test_config_cases\test-20210416173550-lbo-0416cpc", r"C:\Users\JWang294\src\cis\libcalc\auto_report_pipeline1.1\auto_test_config.ini")