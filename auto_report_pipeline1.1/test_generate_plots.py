import pandas, pptx, pptx.util, glob, re, csv, io, PIL
from pptx import Presentation
import cv2
import configparser
import matplotlib.pyplot as plt
import numpy as np


def read_summary_doc(org_dir, list_type):
    i = 0
    iCase = 0
    # doc_list = []
    axis_num_list = []
    # Connect zero point
    axis_num_list.append(0)
    for pic_dir in glob.glob(org_dir + r'\Contact\contact_script_test\summary_doc'):
        with open(pic_dir, "r", encoding="gbk", errors="ignore") as rop_data:
            csv_reader = csv.reader(rop_data)
            for row in csv_reader:
                doc_list = str(row).split()
                #save rop/wob... as a list
                if i >= 2 and i < 3:
                    # doc_list = str(row).split()
                    for type in doc_list:
                        #match type to choose which line to return
                        if (type.find(list_type)) != -1:
                            iCase = doc_list.index(type)
                            break
                #add right num_line to axis_num_list
                elif i > 3:
                    axis_num_list.append(float(doc_list[iCase]))
                i += 1
        break
    return axis_num_list

def generate_contact_plots(org_dir1, org_dir2, config_path):
    cf = configparser.ConfigParser()
    cf.read(config_path, encoding='utf-8')
    template_path = cf.get("merge_template", "dirct")
    prs = Presentation(template_path)
    org_import_dir = (org_dir1, org_dir2)
    time_title1 = cf.get("profile_title", "title1")
    time_title2 = cf.get("profile_title", "title2")
    doc_type_list = ("doc", "tq", "rop")
    pic_left = [int(prs.slide_width*0.08), int(prs.slide_width*0.58)]
    pic_top = int(prs.slide_height*0.3)
    pic_width = int(prs.slide_width*0.4)
    pic_height = int(prs.slide_height*0.6)
    for type in doc_type_list:
        # plt.xticks(rotation=90)
        # plt.gcf().autofmt_xdate()
        # my_x_ticks = np.arange(10000,60000,10000)
        # plt.xticks(my_x_ticks)
        # line 1 points
        x1 = read_summary_doc(org_dir1, "wob")
        y1 = read_summary_doc(org_dir1, type)
        # xpoints = np.array(x1)
        # ypoints = np.array(y1)
        # plotting the line 1 points 
        plt.plot(x1, y1, label = "build" + time_title1, color='c')
        # line 2 points
        x2 = read_summary_doc(org_dir2, "wob")
        y2 = read_summary_doc(org_dir2, type)
        # plotting the line 2 points 
        plt.plot(x2, y2, label = "build" + time_title2, color='m')
        # add Horizontal line
        plt.grid(axis="y")
        # naming the x&y axis
        if type == "doc":
            y_axis_label = "Depth of Cut (in/rev)"
            plot_title = "Depth of Cut-Wellington Shale @ 3ksi"
            # add Horizontal line
            plt.axhline(y2[3], color='orange', lw=2, alpha=0.7)
            # fill with color
            plt.axhspan(0, y2[3], alpha=0.5, color='y')
            # plt.fill_between(x = , y = y2[3], y = 0, step="pre", color = "g", alpha = 0.3)
            plt.ylim(0,0.6)
        elif type == "tq":
            y_axis_label = "Torque (ft-lbf)"
            plot_title = "Torque-Wellington Shale @ 3ksi" 
            plt.axvline(x2[4], color='gold', lw=2, alpha=0.7)
            plt.ylim(0,10000)
            # plt.hlines(0, xmin=0, xmax=10000, color='r', linestyle='-')
            # plt.axis([10000,60000, 2000,10000])
        else:
            y_axis_label = "ROP (ft/hr)"
            plot_title = "ROP-Wellington Shale @ 3ksi" 
            plt.axvline(x2[4], color='orange', lw=2, alpha=0.7)
            plt.ylim(0,450)
            # plt.axis([10000,60000, 50,450])
        plt.xlim(0,)
        plt.xlabel('WOB(lbf)')
        plt.ylabel(y_axis_label)
        # giving a title to my graph
        plt.title(plot_title)
        plt.legend()
        buffer_png = io.BytesIO()
        plt.savefig(buffer_png, format = "png")
        plt.show()
        # # function to show the plot
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        pic1 = slide.shapes.add_picture(buffer_png, pic_left[0], pic_top, pic_width, pic_height)
        plt.close()
        buffer_png.close()
    prs.save(template_path)

if __name__ == "__main__":
    generate_contact_plots(r"C:\Users\JWang294\Documents\projectfile\test_config_cases\test-20210416173550-lbo-0416cpc", r"C:\Users\JWang294\Documents\projectfile\test_config_cases\test-20210416173550-lbo-0416cpc", r"C:\Users\JWang294\src\cis\libcalc\auto_report_pipeline1.1\auto_test_config.ini")  
